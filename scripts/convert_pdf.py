#!/usr/bin/env python3
"""Conversión PDF -> Office con calidad.
- docx: texto editable real (PyMuPDF) + imágenes embebidas, manteniendo estructura por bloques
- pptx: cada página del PDF como imagen en una diapositiva (se ve idéntico)
- xlsx: extracción de texto con PyMuPDF -> openpyxl
Uso: convert_pdf.py <input.pdf> <target> <output>
"""
import sys, os, tempfile, subprocess, re, io

def pdf_to_docx(inp, out):
    """Convierte el PDF a Word con TEXTO EDITABLE real + imágenes embebidas.
    - Extrae bloques de texto ordenados por posición (readable) -> párrafos editables.
    - Extrae las imágenes del PDF y las inserta en el documento.
    """
    import fitz  # PyMuPDF
    from docx import Document
    from docx.shared import Inches, Pt, Emu
    from PIL import Image as PILImage

    pdf = fitz.open(inp)
    d = Document()
    style = d.styles["Normal"]
    style.font.size = Pt(11)

    for page in pdf:
        # --- Texto editable: bloques ordenados por posición (de arriba a abajo) ---
        blocks = page.get_text("blocks")  # [x0,y0,x1,y1,text,block_no,type]
        # type 0 = texto, type 1 = imagen
        blocks_sorted = sorted([b for b in blocks if b[4].strip()], key=lambda b: (round(b[1]), round(b[0])))
        for b in blocks_sorted:
            if b[6] == 0:  # bloque de texto
                text = b[4].strip()
                if text:
                    for line in text.split("\n"):
                        line = line.strip()
                        if line:
                            d.add_paragraph(line)
            elif b[6] == 1:  # bloque de imagen
                try:
                    img_bytes = page.extract_image(b[5])["image"]
                    stream = io.BytesIO(img_bytes)
                    im = PILImage.open(stream)
                    # ancho razonable (máx 5.5in) manteniendo proporción
                    max_w = 5.5
                    w_in = min(max_w, im.width / 96.0)
                    h_in = im.height / im.width * w_in
                    p = d.add_paragraph()
                    p.alignment = 1  # centro
                    p.add_run().add_picture(io.BytesIO(img_bytes), width=Inches(w_in))
                except Exception:
                    pass
        # Salto de página entre páginas
        if page.number < len(pdf) - 1:
            d.add_page_break()
    d.save(out)

def pdf_to_pptx(inp, out):
    from pptx import Presentation
    from pptx.util import Inches
    from PIL import Image
    wd = tempfile.mkdtemp()
    subprocess.run(["pdftoppm", "-png", "-r", "110", inp, os.path.join(wd, "page")], check=True)
    pages = sorted([f for f in os.listdir(wd) if f.endswith(".png")])
    if not pages:
        raise RuntimeError("no se pudo rasterizar el PDF")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for p in pages:
        img_path = os.path.join(wd, p)
        im = Image.open(img_path)
        w, h = im.size
        max_w, max_h = 13.0, 7.0
        ratio = min(max_w / (w / 96.0), max_h / (h / 96.0))
        disp_w = (w / 96.0) * ratio
        disp_h = (h / 96.0) * ratio
        left = (13.333 - disp_w) / 2
        top = (7.5 - disp_h) / 2
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(img_path, Inches(left), Inches(top), Inches(disp_w), Inches(disp_h))
    prs.save(out)

def pdf_to_xlsx(inp, out):
    import fitz  # PyMuPDF
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "PDF"
    row = 1
    doc = fitz.open(inp)
    for page in doc:
        text = page.get_text("text") or ""
        for line in text.split("\n"):
            cells = [c.strip() for c in re.split(r"\s{2,}|\t", line) if c.strip()]
            if not cells:
                continue
            for col, val in enumerate(cells, start=1):
                ws.cell(row=row, column=col, value=val)
            row += 1
    wb.save(out)

def main():
    inp = sys.argv[1]
    target = sys.argv[2]
    out = sys.argv[3]
    if target == "docx":
        pdf_to_docx(inp, out)
    elif target == "pptx":
        pdf_to_pptx(inp, out)
    elif target == "xlsx":
        pdf_to_xlsx(inp, out)
    else:
        raise RuntimeError("target no válido")
    if not os.path.exists(out) or os.path.getsize(out) == 0:
        raise RuntimeError("conversión falló: sin archivo de salida")

if __name__ == "__main__":
    main()
